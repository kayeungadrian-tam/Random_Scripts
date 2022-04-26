import docx2txt
import re
from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import json
from modules.bible_verse import OnlineBibleVerse

def book_handle(book):    
    with open("./assets/books_dict.json", 'r', encoding='utf-8') as f:
        books_dict = json.load(f)
    data = books_dict[book]
    return data



def check_int(str):
        return any(chr.isdigit() for chr in str)


title_img = "./assets/2022 John.png"
background_img = "./assets/2022 John text.png"


params = {
    "title_img": title_img,
    "background_img": background_img,
    "slide_width": Cm(25.5),
    "slide_height": Cm(14.3),
    "font_size": Pt(20),
    "font_color": RGBColor(255, 255, 255),
    "font_name": "Arial",
    "pic_left": Cm(0),
    "pic_top": Cm(0),
    "text_margin": Cm(2),
    "text_center": PP_ALIGN.CENTER
}

class PowerPoint():
    def __init__(self, **kwargs):
        self.prs = Presentation()
        self._setup(**kwargs)
        self.__set_wh()
        self.__set_title()

    def _setup(self, **kwargs):
        print("Setting up...")
        self.width = kwargs.get("slide_width")
        self.height = kwargs.get("slide_height")
        self.font_size = kwargs.get("font_size")
        self.font_name = kwargs.get("font_name")
        self.font_color = kwargs.get("font_color")
        self.title_img = kwargs.get("title_img")
        self.background_img = kwargs.get("background_img")
        self.pic_left = kwargs.get("pic_left")
        self.pic_top = kwargs.get("pic_top")
        self.text_margin = kwargs.get("text_margin")
        self.text_center = kwargs.get("text_center")

    def __set_wh(self):
        self.prs.slide_width = self.width
        self.prs.slide_height = self.height

    def __set_title(self):
        title_slide = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(title_slide)
        pic = slide.shapes.add_picture(self.title_img, self.pic_left, self.pic_top, self.width, self.height)

    def save(self, fname):
        self.prs.save(fname)

    def _set_textboxes_point(self, slide):
        point_tb = slide.shapes.add_textbox(
            self.pic_left + self.text_margin, 
            self.pic_top + (1/3)*self.height, 
            self.width - 2 * self.text_margin, 
            (self.height - 2 * self.text_margin)//2
            )

        point_p = self.__clean_textbox_point(point_tb)
        return point_p

    def add_point_slide(self, ja, en):        
        verse_layout = self.prs.slide_layouts[6]
        text_slide = self.prs.slides.add_slide(verse_layout)
        pic2 = text_slide.shapes.add_picture(self.background_img, self.pic_left, self.pic_top, self.width, self.height)
        tb1_p = self._set_textboxes_point(text_slide)
        tb1_p.text = f"{ja}\n\n{en}"

    def add_img_slide(self, text):
        self.add_point_slide(text, text)

    def __clean_textbox_point(self, tbox):
        tb1 = tbox.text_frame
        tb1.word_wrap = True

        p = tb1.add_paragraph()
        p.font.color.rgb = self.font_color
        p.font.size = self.font_size
        p.font.name = self.font_name
        p.alignment = self.text_center
        p.font.bold = True
        return p



slides = PowerPoint(**params)
BibleVerse = OnlineBibleVerse()


doc_path = "./demo.docx"

my_text = docx2txt.process(doc_path)
# print(my_text)

pattern = r"\[(.*?)\]"
r = re.findall(pattern, my_text)


for idx, key in enumerate(r):
    splitted = key.split(":")
    
    if len(splitted) == 2:
        if splitted[0] == "PIC":
            slides.add_img_slide(f"PIC: {splitted[1]}")
        elif splitted[0] == "PNT":
            slides.add_point_slide(f"Point: {splitted[1]}", f"Point: {splitted[1]}")
    if len(splitted) == 3:
        book = splitted[1].split(" ")[1]
        chapter = int(splitted[1].split(" ")[2])
        verses = splitted[2].split("-")

        data = book_handle(book)
        book_id = int(data["idx"])
        english = ""
        japanese = ""
        # for verse in range(int(verses[0]), int(verses[1]) + 1):
        for verse in verses:
            verse = int(verse)
            # get ja and en
            print("Logging in...")
            BibleVerse.ja_login("./assets/credentials.json", "https://bible.prsi.org/ja/Account/Login")
            print("Logged in.")

            ja_text = BibleVerse.get_verse_ja(book_id, chapter, verse)
            japanese += ja_text + " "

            en_text = BibleVerse.get_verse_en(book, data["code"], chapter, verse)
            english += en_text + " "
        slides.add_point_slide(japanese, english)


        #     # slides.add_img_slide(f"{book} {chapter}:{verse}")
        #     print(book, chapter, verse)

slides.save("demo.pptx")

