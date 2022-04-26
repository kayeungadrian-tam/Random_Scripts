from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from datetime import datetime

title_img = "./assets/2022 John.png"
background_img = "./assets/2022 John text.png"

params = {
    "title_img": title_img,
    "background_img": background_img,
    "slide_width": Cm(25.5),
    "slide_height": Cm(14.3),
    "font_size": Pt(20),
    "verse_size": Pt(18),
    "font_color": RGBColor(255, 255, 255),
    "font_name": "Arial",
    "pic_left": Cm(0),
    "pic_top": Cm(0),
    "text_margin": Cm(2),
    "text_center": PP_ALIGN.CENTER,
    "text_left": PP_ALIGN.LEFT,
    "text_right": PP_ALIGN.RIGHT
    
}

class PowerPoint():
    def __init__(self, **kwargs):
        self.prs = Presentation()
        self._setup(**kwargs)
        self.__set_wh()
        self.__set_title()

    def _setup(self, **kwargs):
        print("Setting up...")
        self.width = kwargs.get("slide_width")
        self.height = kwargs.get("slide_height")
        self.font_size = kwargs.get("font_size")
        self.font_name = kwargs.get("font_name")
        self.font_color = kwargs.get("font_color")
        self.title_img = kwargs.get("title_img")
        self.background_img = kwargs.get("background_img")
        self.pic_left = kwargs.get("pic_left")
        self.pic_top = kwargs.get("pic_top")
        self.text_margin = kwargs.get("text_margin")
        self.text_center = kwargs.get("text_center")
        self.verse_size = kwargs.get("verse_size")
        self.text_left = kwargs.get("text_left")
        self.text_right = kwargs.get("text_right")

    def __set_wh(self):
        self.prs.slide_width = self.width
        self.prs.slide_height = self.height

    def __set_title(self):
        title_slide = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(title_slide)
        pic = slide.shapes.add_picture(self.title_img, self.pic_left, self.pic_top, self.width, self.height)

    def save(self):
        

        time_now = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        self.prs.save(f"./output/{time_now}.pptx")

    def _set_textboxes_point(self, slide):
        point_tb = slide.shapes.add_textbox(
            self.pic_left + self.text_margin, 
            self.pic_top + (1/3)*self.height, 
            self.width - 2 * self.text_margin, 
            (self.height - 2 * self.text_margin)//2
            )

        point_p = self.__clean_textbox_point(point_tb)
        return point_p

    def add_point_slide(self, ja, en):        
        verse_layout = self.prs.slide_layouts[6]
        text_slide = self.prs.slides.add_slide(verse_layout)
        pic2 = text_slide.shapes.add_picture(self.background_img, self.pic_left, self.pic_top, self.width, self.height)
        tb1_p = self._set_textboxes_point(text_slide)
        tb1_p.text = f"{ja}\n\n{en}"

    def add_verse_slides(self, ja, en, ja_ref, en_ref):
        verse_layout = self.prs.slide_layouts[6]
        text_slide = self.prs.slides.add_slide(verse_layout)
        pic2 = text_slide.shapes.add_picture(self.background_img, self.pic_left, self.pic_top, self.width, self.height)
        tb1_p, r1, tb2_p, r2 = self._set_textboxes_verse(text_slide)

        tb1_p.text = f"{ja}"
        r1.text = f"{ja_ref}"
        tb2_p.text = f"{en}"
        r2.text = f"{en_ref}"

    def add_img_slide(self, text):
        self.add_point_slide(text, text)

    def __clean_textbox_point(self, tbox):
        tb1 = tbox.text_frame
        tb1.word_wrap = True

        p = tb1.add_paragraph()
        p.font.color.rgb = self.font_color
        p.font.size = self.font_size
        p.font.name = self.font_name
        p.alignment = self.text_center
        p.font.bold = True
        return p

    def __clean_textbox_verse(self, tbox):
        tb1 = tbox.text_frame
        tb1.word_wrap = True

        p = tb1.add_paragraph()
        p.font.color.rgb = self.font_color
        p.font.size = self.verse_size
        p.font.name = self.font_name
        p.alignment = self.text_left
        p.font.bold = False
        
        p2 = tb1.add_paragraph()
        p2.font.color.rgb = self.font_color
        p2.font.size = self.verse_size
        p2.font.name = self.font_name
        p2.alignment = self.text_right
        
        return p, p2

    def _add_verse(self, tbox, verse):
        p2 = tbox.add_paragraph()
        p2.font.color.rgb = self.font_color
        p2.font.size = self.verse_size
        p2.font.name = self.font_name
        p2.alignment = self.text_right
        return p

    def _set_textboxes_verse(self, slide):
        verse_tb1 = slide.shapes.add_textbox(
            self.pic_left + self.text_margin, 
            self.pic_top + self.text_margin, 
            self.width - 2 * self.text_margin, 
            (self.height - 2 * self.text_margin)//2
            )

        verse_tb2 = slide.shapes.add_textbox(
            self.pic_left + self.text_margin, 
            self.pic_top + (1/2)*self.height, 
            self.width - 2 * self.text_margin, 
            (self.height - 2 * self.text_margin)//2
            )

        verse_p1, ref1 = self.__clean_textbox_verse(verse_tb1)
        
        verse_p2, ref2 = self.__clean_textbox_verse(verse_tb2)

        return verse_p1, ref1, verse_p2, ref2

def test():

    japanese_text = '''
    そこで、ペテロは彼らに言った。「それぞれ罪を赦していただくために、悔い改めて、イエス・キリストの名によってバプテスマを受けなさい。そうすれば、賜物として聖霊を受けます。」
    '''

    ja_ref = "ヨハネの福音書 3:16"

    english_text = '''
    [38] And Peter said to them, “Repent and be baptized every one of you in the name of Jesus Christ for the forgiveness of your sins, and you will receive the gift of the Holy Spirit.”
    '''

    en_ref = "John 3:16"


    time_now = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

    slides = PowerPoint(**params)
    slides.add_point_slide("あいうえお", "Point")
    slides.add_img_slide("Image")
    slides.add_verse_slides(japanese_text, english_text, ja_ref, en_ref)
    slides.save(f"./output/{time_now}.pptx")

if __name__ == "__main__":
    test()